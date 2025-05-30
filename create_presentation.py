#!/usr/bin/env python3
"""
AIエージェントとMCP（Model Context Protocol）に関するプレゼンテーション作成スクリプト
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import requests
from io import BytesIO

def add_image_from_url(slide, url, left, top, width, height):
    """URLから画像をダウンロードしてスライドに追加"""
    try:
        response = requests.get(url)
        if response.status_code == 200:
            image_stream = BytesIO(response.content)
            slide.shapes.add_picture(image_stream, left, top, width, height)
        else:
            # 画像の代わりにテキストボックスでURLを表示
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = f"画像: {url}"
    except Exception as e:
        # エラーの場合もテキストボックスでURLを表示
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = f"画像: {url}"

def create_ai_agents_mcp_presentation():
    """AIエージェントとMCPのプレゼンテーションを作成"""
    
    # プレゼンテーションを作成
    prs = Presentation()
    
    # タイトルスライド
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AIエージェントとMCP\n(Model Context Protocol)"
    subtitle.text = "開発者向けセミナー\n2025年6月"
    
    # スライド2: AIエージェントとは
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "AIエージェントとは？"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "AIエージェントは、ユーザーの代わりに行動し、タスクを実行するAIシステムです"
    
    p = tf.add_paragraph()
    p.text = "主要な例："
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "GitHub Copilot Agent Mode (Visual Studio Code)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "類似ツール: Cline、Cursor"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "GitHub Copilot Coding Agent"
    p.level = 1
    
    # スライド3: AIエージェントエコシステム
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "AIエージェントエコシステム"
    
    # エコシステム図を追加
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    text_frame = textbox.text_frame
    text_frame.text = """AIエージェントエコシステム:
    
                    [Human User]
                         ↓
                   [AI Assistant]
                    ↙    ↓    ↘
            [Code Agent] [Data Agent] [UI Agent]
                 ↓           ↓           ↓
              [GitHub]   [Database]   [Browser]
                 ↓           ↓           ↓
            [Repository] [Analytics] [Web Apps]
    
    全てのエージェントがMCPプロトコルで相互接続"""
    
    # スライド4: GitHub Copilot Coding Agent
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "GitHub Copilot Coding Agent"
    
    # GitHub Copilot Coding Agentの画像を追加
    add_image_from_url(slide, 
                      "https://www.publickey1.jp/2025/github-codingagent-01.png",
                      Inches(1), Inches(2), Inches(8), Inches(4))
    
    # 画像の下にソース情報を追加
    textbox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    text_frame = textbox.text_frame
    text_frame.text = "出典: GitHub Copilot Coding Agent発表資料"
    
    # スライド5: MCPとは
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "MCP (Model Context Protocol) とは？"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "AIエージェント同士が連携するためのプロトコル"
    
    p = tf.add_paragraph()
    p.text = "Microsoftが全面採用を発表"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Windows、Microsoft 365などの主要製品でサポート"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "AIエージェント間のスムーズな情報交換と連携を実現"
    p.level = 1
    
    # スライド6: MCP概念図
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "MCP概念図"
    
    # MCP概念図を追加
    add_image_from_url(slide,
                      "https://www.publickey1.jp/2025/windows-support-mcp.png",
                      Inches(1), Inches(1.5), Inches(8), Inches(5))
    
    # ソース情報を追加
    textbox = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.8))
    text_frame = textbox.text_frame
    text_frame.text = "出典: [速報] MicrosoftがWindowsでMCPサポートを発表、AIエージェントがWindowsやアプリと連携可能に - Publickey\nhttps://www.publickey1.jp/blog/25/windowsmcpaiwindows.html"
    
    # スライド7: Microsoft BUILD 2025
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Microsoft BUILD 2025"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "2025年5月20-23日開催"
    
    p = tf.add_paragraph()
    p.text = "テーマ: 「AIエージェントの時代」"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "CEO Nadellaが従来のAIモデルからAIエージェントへの進化を強調"
    p.level = 1
    
    # Satya Nadellaの画像を追加
    add_image_from_url(slide,
                      "https://image.itmedia.co.jp/news/articles/2505/20/l_yu_satya.jpg",
                      Inches(1), Inches(3.5), Inches(8), Inches(3))
    
    # スライド8: 主な発表内容
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "主な発表内容"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "MCPサーバーレジストリの提供"
    
    p = tf.add_paragraph()
    p.text = "WindowsでのMCPサポート"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Microsoft 365でのMCPサポート"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "GitHub Copilot Coding Agentの発表"
    p.level = 0
    
    # スライド9: マルチエージェントオーケストレーション
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "マルチエージェントオーケストレーション"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "複数のAIエージェントが連携して複雑なタスクを実行"
    
    p = tf.add_paragraph()
    p.text = "各エージェントが専門的な役割を担当"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "MCPプロトコルにより効率的な情報共有"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "サービス、OS、アプリケーション、Webページとの連携"
    p.level = 1
    
    # 概念図を追加（シンプルなテキストベースの図）
    textbox = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2.5))
    text_frame = textbox.text_frame
    text_frame.text = """マルチエージェント構成例:
    
    [UI Agent] ←→ [MCP Protocol] ←→ [API Agent]
         ↕                              ↕
    [User Interface]              [External Services]
         ↕                              ↕  
    [Analytics Agent] ←→ [MCP] ←→ [Database Agent]"""
    
    # スライド10: MCPサーバー
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "MCPサーバー"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "AIエージェントと外部システム間の橋渡し役"
    
    p = tf.add_paragraph()
    p.text = "様々なサービスやAPIとの接続を提供"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "標準化されたインターフェースで一貫性を保証"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Claude Desktopでの利用デモも提供"
    p.level = 1
    
    # MCPサーバーアーキテクチャ図を追加
    textbox = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2.5))
    text_frame = textbox.text_frame
    text_frame.text = """MCPサーバーアーキテクチャ:
    
    [AI Agent] ←→ [MCP Protocol] ←→ [MCP Server]
                                          ↓
                                   [Adapter Layer]
                                          ↓
                        [Database] [API] [File System] [Web Services]"""
    
    # スライド11: MCPサーバー実装例
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "MCPサーバー実装例"
    
    # 実装例の写真/スクリーンショットとして、Claude Desktopの例を追加
    add_image_from_url(slide,
                      "https://docs.anthropic.com/en/docs/_next/image?url=%2Fen%2Fdocs%2F_next%2Fstatic%2Fmedia%2Fmcp-architecture.bfaa5e8c.png&w=1080&q=75",
                      Inches(0.5), Inches(2), Inches(9), Inches(4.5))
    
    # ソース情報を追加
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.8))
    text_frame = textbox.text_frame
    text_frame.text = "出典: Anthropic MCP Documentation - Claude Desktop MCP Integration Example"
    
    # スライド12: .NETでのMCPサーバー開発
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = ".NETでのMCPサーバー開発"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Microsoft環境での開発支援"
    
    p = tf.add_paragraph()
    p.text = ".NET Framework/Coreでの実装"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "豊富なライブラリとツールチェーンの活用"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "実践的なデモンストレーション付き"
    p.level = 1
    
    # スライド13: Visual Studio Codeでの利用
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Visual Studio Codeでの利用"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "開発者にとって身近な環境での実装"
    
    p = tf.add_paragraph()
    p.text = "拡張機能によるMCPサポート"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "AIエージェントとの直接的な連携"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "実際のワークフローでのデモンストレーション"
    p.level = 1
    
    # スライド14: 未来への展望
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "未来への展望"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "アプリケーションとサービスの性質が大きく変化"
    
    p = tf.add_paragraph()
    p.text = "近年提案されているマイクロアーキテクチャと同様の進化"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "HTTP/HTTPSでマイクロサービスを接続するように"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "MCPでAIエージェントを接続"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Microsoft、その他多くのIT企業がMCPを採用"
    p.level = 1
    
    # スライド15: まとめ
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "まとめ"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "AIエージェントの時代が到来"
    
    p = tf.add_paragraph()
    p.text = "MCPがAIエージェント連携の標準プロトコルに"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "開発者は新しいアーキテクチャパターンへの対応が必要"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "実践的なデモと開発環境の整備が進行中"
    p.level = 1
    
    # スライド16: 参考資料
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "参考資料"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Microsoft Build 2025 公式情報"
    
    p = tf.add_paragraph()
    p.text = "https://www.itmedia.co.jp/news/articles/2505/20/news097.html"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "https://www.publickey1.jp/blog/25/windowsmcpaiwindows.html"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "GitHub Copilot 関連"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "https://github.com/features/copilot"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "https://docs.github.com/en/copilot"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "MCP (Model Context Protocol) 関連"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "https://modelcontextprotocol.org/"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "https://github.com/modelcontextprotocol"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "AI エージェント関連"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "https://openai.com/research/agents"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "https://claude.ai/"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "https://cursor.sh/"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "https://github.com/cline/cline"
    p.level = 1
    
    # プレゼンテーションを保存
    output_path = "/home/runner/work/2025.06.AIAgentsSeminarSlide/2025.06.AIAgentsSeminarSlide/AIエージェントとMCP_セミナー.pptx"
    prs.save(output_path)
    print(f"プレゼンテーションが作成されました: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_ai_agents_mcp_presentation()